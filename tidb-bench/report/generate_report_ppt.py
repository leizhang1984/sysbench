from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
OUTPUT = BASE_DIR / "tidb-dsv5-vs-dsv6-summary.pptx"
IMG_DIR = BASE_DIR / "images"


def set_slide_bg(slide, color_rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color_rgb)


def add_title(slide, text, left=0.6, top=0.3, width=12.1, height=0.7, color=(255, 255, 255)):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*color)
    p.alignment = PP_ALIGN.LEFT


def add_subtitle(slide, text, left=0.6, top=1.15, width=12.1, height=0.5, color=(220, 232, 255)):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(*color)
    p.alignment = PP_ALIGN.LEFT


def add_bullets(slide, items, left, top, width, height, font_size=20, color=(20, 20, 20)):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(6)
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(*color)


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(6.85), Inches(12.0), Inches(0.35))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Calibri"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(110, 110, 110)
    p.alignment = PP_ALIGN.RIGHT


def add_tag(slide, text, left, top, width=2.2, height=0.45):
    shape = slide.shapes.add_shape(
        autoshape_type_id=1,
        left=Inches(left),
        top=Inches(top),
        width=Inches(width),
        height=Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(23, 49, 102)
    shape.line.fill.background()

    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)


def add_arch_box(slide, title, lines, left, top, width, height, fill_rgb, title_rgb=(255, 255, 255)):
    container = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    container.fill.solid()
    container.fill.fore_color.rgb = RGBColor(*fill_rgb)
    container.line.color.rgb = RGBColor(210, 220, 240)

    tf = container.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.name = "Microsoft YaHei"
    p_title.font.size = Pt(12)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(*title_rgb)
    p_title.alignment = PP_ALIGN.CENTER

    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.name = "Calibri"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(*title_rgb)
        p.alignment = PP_ALIGN.CENTER


def add_icon_badge(slide, label, left, top, size=0.34, fill_rgb=(255, 255, 255), text_rgb=(33, 33, 33)):
    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(left),
        Inches(top),
        Inches(size),
        Inches(size),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(*fill_rgb)
    badge.line.fill.background()
    tf = badge.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = label
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Calibri"
    p.font.bold = True
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(*text_rgb)


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, (12, 32, 74))
    add_title(slide, "TiDB DSv5 vs DSv6 压测复盘汇报", color=(255, 255, 255))
    add_subtitle(slide, "同资源规格集群对比：吞吐、延迟与系统资源利用 | 2026-06", color=(192, 213, 255))
    add_bullets(
        slide,
        [
            "测试结论：DSv6 在全部场景性能领先，适合作为推荐机型",
            "核心收益：QPS 平均 +9.9%，延迟平均下降 7%~11%",
            "重点发现：中高并发（100/200）下优势更明显",
        ],
        left=0.8,
        top=2.0,
        width=11.5,
        height=2.8,
        font_size=24,
        color=(236, 242, 255),
    )
    add_tag(slide, "汇报摘要", 0.8, 5.0)
    add_footer(slide, "数据来源：tidb-bench/report/tidb-dsv5-vs-dsv6-report.md")

    # Slide 2: Architecture (clean side-by-side comparison)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, (245, 248, 255))
    add_title(slide, "部署架构图：两套独立集群并行对照", color=(15, 35, 75))

    # Left cluster: DSv5
    cluster5 = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.65),
        Inches(1.25),
        Inches(5.95),
        Inches(5.35),
    )
    cluster5.fill.solid()
    cluster5.fill.fore_color.rgb = RGBColor(229, 238, 255)
    cluster5.line.color.rgb = RGBColor(99, 130, 190)
    tf5 = cluster5.text_frame
    tf5.clear()
    tf5.paragraphs[0].text = "集群 A：DSv5 (Standard_D8s_v5)"
    tf5.paragraphs[0].font.name = "Microsoft YaHei"
    tf5.paragraphs[0].font.size = Pt(16)
    tf5.paragraphs[0].font.bold = True
    tf5.paragraphs[0].font.color.rgb = RGBColor(27, 53, 101)

    add_arch_box(slide, "客户端", ["clientvm01", "D32s_v6"], 1.00, 2.0, 1.55, 1.0, (66, 109, 184))
    add_arch_box(slide, "LB", ["10.142.0.10:4000"], 2.85, 2.0, 1.9, 1.0, (66, 109, 184))
    add_arch_box(slide, "TiDB + PD", ["3 节点", "10.142.0.11/12/13"], 5.00, 1.8, 1.35, 1.4, (66, 109, 184))
    add_arch_box(slide, "TiKV", ["3 节点", "10.142.0.21/22/23", "AZ1/AZ2/AZ3"], 5.00, 3.5, 1.35, 1.7, (76, 124, 206))

    flow_a1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(2.55), Inches(2.35), Inches(0.30), Inches(0.24))
    flow_a1.fill.solid()
    flow_a1.fill.fore_color.rgb = RGBColor(58, 92, 163)
    flow_a1.line.fill.background()
    flow_a2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(4.75), Inches(2.35), Inches(0.25), Inches(0.24))
    flow_a2.fill.solid()
    flow_a2.fill.fore_color.rgb = RGBColor(58, 92, 163)
    flow_a2.line.fill.background()

    # Right cluster: DSv6
    cluster6 = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(6.75),
        Inches(1.25),
        Inches(5.95),
        Inches(5.35),
    )
    cluster6.fill.solid()
    cluster6.fill.fore_color.rgb = RGBColor(229, 246, 236)
    cluster6.line.color.rgb = RGBColor(74, 151, 109)
    tf6 = cluster6.text_frame
    tf6.clear()
    tf6.paragraphs[0].text = "集群 B：DSv6 (Standard_D8s_v6)"
    tf6.paragraphs[0].font.name = "Microsoft YaHei"
    tf6.paragraphs[0].font.size = Pt(16)
    tf6.paragraphs[0].font.bold = True
    tf6.paragraphs[0].font.color.rgb = RGBColor(22, 81, 49)

    add_arch_box(slide, "客户端", ["clientvm02", "D32s_v6"], 7.10, 2.0, 1.55, 1.0, (44, 133, 88))
    add_arch_box(slide, "LB", ["10.142.0.30:4000"], 8.95, 2.0, 1.9, 1.0, (44, 133, 88))
    add_arch_box(slide, "TiDB + PD", ["3 节点", "10.142.0.31/32/33"], 11.10, 1.8, 1.35, 1.4, (44, 133, 88))
    add_arch_box(slide, "TiKV", ["3 节点", "10.142.0.41/42/43", "AZ1/AZ2/AZ3"], 11.10, 3.5, 1.35, 1.7, (57, 152, 99))

    flow_b1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(8.65), Inches(2.35), Inches(0.30), Inches(0.24))
    flow_b1.fill.solid()
    flow_b1.fill.fore_color.rgb = RGBColor(66, 147, 100)
    flow_b1.line.fill.background()
    flow_b2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(10.85), Inches(2.35), Inches(0.25), Inches(0.24))
    flow_b2.fill.solid()
    flow_b2.fill.fore_color.rgb = RGBColor(66, 147, 100)
    flow_b2.line.fill.background()

    add_tag(slide, "配置差异说明", 0.8, 6.0, width=2.2, height=0.38)
    add_bullets(
        slide,
        [
            "DSv5（集群A）：VM=Standard_D8s_v5；OS=CentOS 7.9",
            "DSv6（集群B）：VM=Standard_D8s_v6；OS=Rocky Linux 9.8",
            "共同项：TiDB v8.5.6、节点拓扑(3xTiDB+PD + 3xTiKV)、磁盘规格一致",
        ],
        left=0.8,
        top=6.35,
        width=12.0,
        height=0.6,
        font_size=12,
        color=(55, 55, 55),
    )
    add_footer(slide, "左：DSv5 + CentOS 7.9 | 右：DSv6 + Rocky Linux 9.8")

    # Slide 3: Core performance results
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, (255, 255, 255))
    add_title(slide, "核心性能结果：DSv6 全面领先", color=(15, 35, 75))

    img_qps = IMG_DIR / "qps_summary.png"
    if img_qps.exists():
        slide.shapes.add_picture(str(img_qps), Inches(0.7), Inches(1.2), height=Inches(4.3))

    add_tag(slide, "关键数字", 8.0, 1.3, width=4.5)
    add_bullets(
        slide,
        [
            "QPS 提升范围：+7.6% ~ +12.7%",
            "QPS 平均提升：+9.9%",
            "平均延迟下降：7% ~ 11%",
            "最显著点：100 并发（RO +11.04%，RW +12.71%）",
            "200 并发峰值：RO 71817 QPS，RW 63565 QPS",
        ],
        left=8.05,
        top=1.9,
        width=4.9,
        height=4.8,
        font_size=16,
        color=(30, 30, 30),
    )
    add_footer(slide, "图表：report/images/qps_summary.png")

    # Slide 4: System metrics view
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, (247, 250, 255))
    add_title(slide, "系统指标侧验证：性能提升来源", color=(15, 35, 75))

    img_ro_tidb = IMG_DIR / "ro_tidb.png"
    img_rw_tidb = IMG_DIR / "rw_tidb.png"
    if img_ro_tidb.exists():
        slide.shapes.add_picture(str(img_ro_tidb), Inches(0.7), Inches(1.3), width=Inches(6.1))
    if img_rw_tidb.exists():
        slide.shapes.add_picture(str(img_rw_tidb), Inches(6.95), Inches(1.3), width=Inches(6.1))

    add_bullets(
        slide,
        [
            "同并发下 DSv6 的 TiDB CPU 利用率更高，说明单位时间处理了更多请求",
            "软中断占比多数场景更低，网络/中断处理开销更小",
            "高并发时 TiDB 接近饱和、TiKV 仍有余量 -> 瓶颈主要在 SQL 计算层",
        ],
        left=0.75,
        top=5.35,
        width=12.0,
        height=1.3,
        font_size=16,
        color=(30, 30, 30),
    )
    add_footer(slide, "图表：report/images/ro_tidb.png, report/images/rw_tidb.png")

    # Slide 5: Final recommendation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, (18, 38, 80))
    add_title(slide, "结论与建议", color=(255, 255, 255))

    add_tag(slide, "结论", 0.8, 1.3)
    add_bullets(
        slide,
        [
            "在 12 个测试场景中，DSv6 全部优于 DSv5（吞吐更高、延迟更低）",
            "收益稳定且在中高并发更突出，适合作为默认升级方向",
            "对当前 OLTP 负载，计算实例升级可直接换取业务性能收益",
        ],
        left=0.85,
        top=1.9,
        width=12.0,
        height=2.1,
        font_size=20,
        color=(235, 241, 255),
    )

    add_tag(slide, "落地建议", 0.8, 4.25)
    add_bullets(
        slide,
        [
            "生产集群优先采用 DSv6 系列并开展分批替换",
            "针对 TiDB 层瓶颈，持续优化 SQL 与连接池参数",
            "补充长稳压（>= 1h）与更高并发（300+）验证容量上限",
        ],
        left=0.85,
        top=4.85,
        width=12.0,
        height=1.8,
        font_size=18,
        color=(235, 241, 255),
    )
    add_footer(slide, "建议：以 DSv6 为标准规格进入后续容量与成本评估")

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    output_file = build_presentation()
    print(f"PPT generated: {output_file}")
