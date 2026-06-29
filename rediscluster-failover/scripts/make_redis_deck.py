#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""Generate a 5-slide reporting deck for the Redis 6 cluster failover report."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# ---- palette (Ocean Gradient, matches RocketMQ deck) ----
CRIMSON = RGBColor(0x06, 0x5A, 0x82)   # deep blue (dominant; name kept for minimal diff)
DEEP    = RGBColor(0x06, 0x5A, 0x82)   # deep blue
MID     = RGBColor(0x21, 0x29, 0x5C)   # midnight (dark bg)
ICE     = RGBColor(0xE8, 0xF1, 0xF5)   # very light
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
INK     = RGBColor(0x21, 0x29, 0x2E)
MUTED   = RGBColor(0x6B, 0x7A, 0x85)
GREEN   = RGBColor(0x1B, 0x9E, 0x77)   # good / 0-loss
AMBER   = RGBColor(0xE8, 0x8A, 0x1A)   # warning / SIGSTOP
BLUE    = RGBColor(0x1C, 0x72, 0x93)   # teal accent
PURPLE  = RGBColor(0x6A, 0x4C, 0x93)

EMU = 914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width  = Emu(int(SW * EMU))
prs.slide_height = Emu(int(SH * EMU))
blank = prs.slide_layouts[6]


def rect(slide, x, y, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    if isinstance(runs[0], tuple):
        runs = [runs]
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line_spacing
        for t, sz, col, bold in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold
            r.font.name = "Calibri"
    return tb


def heading(slide, title, sub=None):
    rect(slide, 0, 0, SW, SH, WHITE)
    rect(slide, 0, 0, SW, 1.12, CRIMSON)
    rect(slide, 0, 1.12, SW, 0.06, AMBER)
    txt(slide, 0.55, 0.12, 12.2, 0.52, [[(title, 22, WHITE, True)]], anchor=MSO_ANCHOR.TOP)
    if sub:
            txt(slide, 0.55, 0.74, 12.2, 0.32, [[(sub, 11, RGBColor(0xCF, 0xE3, 0xEC), False)]])


# ============================================================
# SLIDE 1 — title + key data
# ============================================================
s1 = prs.slides.add_slide(blank)
rect(s1, 0, 0, SW, SH, MID)
rect(s1, 0, 0, 4.55, SH, CRIMSON)
rect(s1, 4.55, 0, 0.06, SH, AMBER)

txt(s1, 0.55, 0.55, 3.6, 0.5, [[("集群故障转移测试", 16, ICE, False)]])
txt(s1, 0.55, 1.0, 3.6, 2.0,
    [[("Redis 6.2", 32, WHITE, True)],
     [("Cluster 故障转移", 28, WHITE, True)],
     [("实测分析", 28, WHITE, True)]],
    line_spacing=1.05)
txt(s1, 0.55, 3.2, 3.6, 2.6,
    [[("测试方法", 13, RGBColor(0x9F,0xC6,0xD8), True)],
     [("• 3 主 3 从集群，16384 槽全覆盖", 12.5, ICE, False)],
     [("• 主从强制跨 3 个可用区", 12.5, ICE, False)],
     [("• JedisCluster 探针持续读写", 12.5, ICE, False)],
     [("• 注入两类故障：kill -9 / SIGSTOP", 12.5, ICE, False)],
     [("• 全程逐秒采集 + 写后回读校验", 12.5, ICE, False)]],
    line_spacing=1.18, space_after=6)
txt(s1, 0.55, 6.6, 3.7, 0.5,
    [[("Redis 6.2.14 · node-timeout=15s · appendfsync everysec", 9.5, RGBColor(0x8F,0xA9,0xB8), False)]])

txt(s1, 4.95, 0.55, 7.9, 0.55,
    [[("核心数据：两类故障均自动转移、零数据丢失", 21, WHITE, True)]])

# baseline strip
rect(s1, 4.95, 1.3, 7.85, 0.92, RGBColor(0x0E, 0x3D, 0x57), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s1, 5.2, 1.42, 1.8, 0.7, [[("基线", 13, RGBColor(0x7F,0xD4,0xF0), True)], [("无故障稳态", 10, RGBColor(0xB8,0xC9,0xD4), False)]])
base = [("2,431", "QPS 吞吐"), ("0.86ms", "P50 延迟"), ("2.60ms", "P99 延迟"), ("0", "失败/丢失")]
bx = 6.95
for val, lab in base:
    txt(s1, bx, 1.4, 1.5, 0.72, [[(val, 20, WHITE, True)], [(lab, 9.5, RGBColor(0xB8,0xC9,0xD4), False)]], line_spacing=0.95)
    bx += 1.48

# two scenario cards
cards = [
    ("场景 A", "kill -9 主进程（崩溃，回 RST）", "≈18.6s", "+17.2s", "FAIL 判定", "全量", "重同步", AMBER),
    ("场景 B", "SIGSTOP 冻结主（静默，无 RST）", "≈23s", "+22.5s", "FAIL 判定", "部分", "重同步", CRIMSON),
]
cx0, cy0 = 4.95, 2.45
cw, ch = 3.82, 2.05
gx = 0.21
for i, (tag, sub, big, fail_t, fail_lab, resync, resync_lab, accent) in enumerate(cards):
    x = cx0 + i * (cw + gx); y = cy0
    rect(s1, x, y, cw, ch, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s1, x, y, 0.12, ch, accent)
    txt(s1, x + 0.28, y + 0.16, cw - 0.4, 0.6, [[(tag, 15, accent, True)], [(sub, 11, INK, True)]], line_spacing=1.0)
    txt(s1, x + 0.28, y + 0.9, 1.85, 0.95,
        [[(big, 30, accent, True)], [("客户端中断", 10.5, MUTED, False)]], line_spacing=0.95)
    txt(s1, x + 2.25, y + 0.92, 1.4, 1.0,
        [[(fail_t, 17, INK, True)], [(fail_lab, 9.5, MUTED, False)],
         [(resync + "重同步", 11.5, INK, True)]], line_spacing=1.05)

# bottom insight bar
rect(s1, 4.95, 4.72, cw*2+gx, 2.18, RGBColor(0x0E, 0x3D, 0x57), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s1, 5.25, 4.9, 7.3, 0.4, [[("关键洞察", 14, RGBColor(0x7F,0xD4,0xF0), True)]])
txt(s1, 5.25, 5.3, 7.35, 1.5,
    [[("• 中断窗口 ≈ cluster-node-timeout（15s）+ 选举(~1s) + 客户端刷新", 12.5, ICE, False)],
     [("• 静默故障（无 RST）比进程崩溃更慢、更“黏”：连续 22s 吞吐绝对归零", 12.5, ICE, False)],
     [("• 单分片故障会拖垮整个客户端吞吐（线程被重试榨干）", 12.5, ICE, False)],
     [("• 写后回读校验：335,870 次成功写，零丢失（前提：复制已追平）", 12.5, ICE, False)]],
    line_spacing=1.22, space_after=5)

# ============================================================
# SLIDE 2 — deployment architecture (3 masters / 3 replicas cross-AZ)
# ============================================================
s2 = prs.slides.add_slide(blank)
heading(s2, "部署架构：3 主 3 从，主从强制跨可用区", "Azure westus3 · 6×D8s_v6 · 每分片主与副本分属不同 AZ，每 AZ 各承载 1 主 + 1 从")

az = [
    ("AZ-1", RGBColor(0xDD,0xEC,0xF7), DEEP,   ("A 主 7001", GREEN, True),  ("C 从 7004", BLUE, False)),
    ("AZ-2", RGBColor(0xD5,0xF0,0xE4), GREEN,  ("B 主 7002", GREEN, True),  ("A 从 7005", BLUE, False)),
    ("AZ-3", RGBColor(0xEC,0xE2,0xF5), PURPLE, ("C 主 7003", GREEN, True),  ("B 从 7006", BLUE, False)),
]
px0, py0 = 0.7, 1.55
pw, ph = 3.85, 3.05
gp = 0.28
for i, (azn, bg, edge, m, r) in enumerate(az):
    x = px0 + i * (pw + gp)
    rect(s2, x, py0, pw, ph, bg, line=edge)
    rect(s2, x, py0, pw, 0.5, edge)
    txt(s2, x + 0.2, py0 + 0.06, pw - 0.4, 0.4, [[(azn + "（可用区）", 14, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    # master node
    mt, mc, _ = m
    rect(s2, x + 0.45, py0 + 0.78, pw - 0.9, 0.82, GREEN, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s2, x + 0.45, py0 + 0.78, pw - 0.9, 0.82, [[(mt, 16, WHITE, True)], [("Master · 可写", 10, WHITE, False)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.95)
    # replica node
    rt, rc, _ = r
    rect(s2, x + 0.45, py0 + 1.85, pw - 0.9, 0.82, BLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s2, x + 0.45, py0 + 1.85, pw - 0.9, 0.82, [[(rt, 16, WHITE, True)], [("Replica · 只读副本", 10, WHITE, False)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.95)

# replication arrows note (A:AZ1->AZ2, B:AZ2->AZ3, C:AZ3->AZ1)
txt(s2, 0.7, 4.78, 11.9, 0.4,
    [[("异步复制（跨 AZ）：", 13, INK, True),
      ("A 主(AZ1) → A 从(AZ2)　B 主(AZ2) → B 从(AZ3)　C 主(AZ3) → C 从(AZ1)", 13, DEEP, True)]])

# bottom facts cards
facts = [
    ("16384 槽全覆盖", "3 分片均分，cluster-require-full-coverage=no：单分片转移期间其余分片继续服务", BLUE),
    ("Gossip 故障判定", "节点经集群总线交换 PING/PONG，PFAIL→FAIL 需多数主确认；node-timeout=15s 决定灵敏度", AMBER),
    ("AZ 级容灾", "任一 AZ 整体故障，每个分片仍有跨 AZ 存活副本可自动升主，保住已复制数据", GREEN),
]
fx0, fy0 = 0.7, 5.4
fw, fh = 3.85, 1.55
for i, (t, b, accent) in enumerate(facts):
    x = fx0 + i * (fw + gp)
    rect(s2, x, fy0, fw, fh, ICE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s2, x, fy0, fw, 0.1, accent)
    txt(s2, x + 0.25, fy0 + 0.22, fw - 0.5, 0.4, [[(t, 13.5, accent, True)]])
    txt(s2, x + 0.25, fy0 + 0.66, fw - 0.5, 0.8, [[(b, 10.5, INK, False)]], line_spacing=1.08)

# ============================================================
# SLIDE 3 — two-scenario comparison (timeline + chart)
# ============================================================
s3 = prs.slides.add_slide(blank)
heading(s3, "两类故障端到端对比：进程崩溃 vs 整机静默", "同一集群、同一探针，唯一变量是故障类型；检测窗口（≈node-timeout）主导中断时长")

# left: comparison chart (outage seconds + failures)
chart_data = CategoryChartData()
chart_data.categories = ["客户端中断(s)", "FAIL判定(s)", "选举耗时(s)"]
chart_data.add_series("场景A kill -9", (18.6, 17.2, 1.02))
chart_data.add_series("场景B SIGSTOP", (23.0, 22.5, 0.64))
gframe = s3.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                             Inches(0.55), Inches(1.35), Inches(6.1), Inches(3.4), chart_data)
ch3 = gframe.chart
ch3.has_legend = True
ch3.legend.position = XL_LEGEND_POSITION.BOTTOM
ch3.legend.include_in_layout = False
ch3.legend.font.size = Pt(10)
plot = ch3.plots[0]
plot.gap_width = 80
ser_colors = [AMBER, CRIMSON]
for s, c in zip(ch3.series, ser_colors):
    s.format.fill.solid(); s.format.fill.fore_color.rgb = c
cat_ax = ch3.category_axis; cat_ax.tick_labels.font.size = Pt(10)
val_ax = ch3.value_axis; val_ax.tick_labels.font.size = Pt(9)
val_ax.has_major_gridlines = True

txt(s3, 0.55, 4.85, 6.1, 0.4,
    [[("选举本身仅 ~1s；中断时长主要被 15s 检测窗口决定。", 11.5, MUTED, False)]])

# right: detail table-like panel
tx, ty, tw = 6.95, 1.35, 5.85
rect(s3, tx, ty, tw, 3.4, ICE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rows = [
    ("维度", "场景A kill -9", "场景B SIGSTOP", True),
    ("副本 TCP 感知", "6ms（收到 RST）", "无即时感知（靠超时）", False),
    ("集群标记 FAIL", "+17.2s", "+22.5s（更黏）", False),
    ("客户端净不可用", "≈18.6s（有涓流）", "≈23s（绝对 0）", False),
    ("恢复瞬间 Max 延迟", "4175 ms", "2354 ms", False),
    ("旧节点回归", "全量重同步", "部分重同步(~2.5MB)", False),
    ("数据丢失 (RPO)", "0", "0", False),
]
colx = [tx + 0.25, tx + 2.05, tx + 4.0]
ry = ty + 0.2
for r0, r1, r2, hd in rows:
    col = CRIMSON if hd else INK
    sz = 11.5 if hd else 11
    bold = hd
    txt(s3, colx[0], ry, 1.8, 0.4, [[(r0, sz, col, bold)]])
    txt(s3, colx[1], ry, 1.95, 0.4, [[(r1, sz, (AMBER if not hd else col), bold)]])
    txt(s3, colx[2], ry, 1.85, 0.4, [[(r2, sz, (CRIMSON if not hd else col), bold)]])
    if hd:
        rect(s3, tx + 0.2, ry + 0.36, tw - 0.4, 0.02, RGBColor(0xC4,0xD6,0xDE))
    ry += 0.445

# bottom takeaways
takes = [
    ("静默故障更致命", "无 RST 时请求石沉大海，线程挂满超时，吞吐连续 22s 绝对归零", CRIMSON),
    ("防脑裂靠 configEpoch", "单调纪元 + 多数派授权，保证同槽同一时刻只有一个合法主", PURPLE),
    ("自动恢复零干预", "旧主重启/解冻后看到更高 epoch 自动降为副本回归", GREEN),
]
kx0, ky0 = 0.55, 5.45
kw, kh = 4.05, 1.5
gk = 0.27
for i, (t, b, accent) in enumerate(takes):
    x = kx0 + i * (kw + gk)
    rect(s3, x, ky0, kw, kh, WHITE, line=RGBColor(0xD6,0xE2,0xE8), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s3, x, ky0, 0.1, kh, accent)
    txt(s3, x + 0.25, ky0 + 0.2, kw - 0.45, 0.4, [[(t, 13.5, accent, True)]])
    txt(s3, x + 0.25, ky0 + 0.64, kw - 0.45, 0.8, [[(b, 10.5, INK, False)]], line_spacing=1.1)

# ============================================================
# SLIDE 4 — data integrity / RPO
# ============================================================
s4 = prs.slides.add_slide(blank)
heading(s4, "数据完整性与 RPO：零丢失是“有条件的”", "Redis Cluster = 最终一致 + 异步复制；主对客户端 ACK 不等副本确认 → RPO 可能 > 0")

# left: safeguards
rect(s4, 0.55, 1.35, 6.0, 5.4, ICE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s4, 0.8, 1.55, 5.6, 0.4, [[("已启用的完整性防线", 15, GREEN, True)]])
safe = [
    ("AOF 持久化", "appendonly + appendfsync everysec，崩溃重启可回放"),
    ("最优副本升主", "选举携 master_repl_offset，偏移最大者优先，最小化丢失"),
    ("复制积压缓冲", "repl-backlog 100mb，解冻后部分重同步（实证）"),
    ("configEpoch 仲裁", "单调纪元 + 多数派，防脑裂写冲突"),
    ("跨 AZ 主从", "单 AZ 故障分片仍存活，保住已复制数据"),
    ("拒绝陈旧写", "旧主见更高 epoch 立即降副本，杜绝分叉"),
]
sy = 2.05
for t, b in safe:
    rect(s4, 0.8, sy + 0.06, 0.14, 0.6, GREEN)
    txt(s4, 1.05, sy, 5.3, 0.7, [[(t, 12.5, INK, True)], [(b, 10.5, MUTED, False)]], line_spacing=1.0)
    sy += 0.77

# right: loss scenarios + why-zero
rect(s4, 6.8, 1.35, 6.0, 2.95, RGBColor(0xFB,0xF1,0xE2), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s4, 7.05, 1.55, 5.5, 0.4, [[("可能丢数据的场景（须生产前评估）", 15, AMBER, True)]])
loss = [
    "异步复制窗口：已 ACK 未复制的写随旧主丢失（最常见）",
    "AOF fsync 间隙：进程崩溃 + 整机断电丢最后 ~1s",
    "双重故障：主与唯一副本同时失联，分片无法转移",
    "脑裂边缘 / full-resync 期间主再故障",
]
ly = 2.05
for s in loss:
    txt(s4, 7.05, ly, 5.5, 0.5, [[("•  ", 12, AMBER, True), (s, 11.5, INK, False)]], line_spacing=1.05)
    ly += 0.55

rect(s4, 6.8, 4.5, 6.0, 2.25, RGBColor(0x1B, 0x4D, 0x3E), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s4, 7.05, 4.68, 5.5, 0.4, [[("本测为何 0 丢失 + 如何做到“实质不丢”", 14, RGBColor(0xA8,0xE6,0xCF), True)]])
txt(s4, 7.05, 5.12, 5.55, 1.6,
    [[("• 写入适中、复制已追平（偏移两线重合），异步窗口≈0", 11.5, ICE, False)],
     [("• 写操作幂等 + 重试退避：覆盖检测窗口，超时即补回", 11.5, ICE, False)],
     [("• min-replicas-to-write / WAIT N：半同步收窄 RPO", 11.5, ICE, False)],
     [("• 对账补偿兜底；务必按峰值写入复测 RPO", 11.5, ICE, False)]],
    line_spacing=1.22, space_after=4)

# ============================================================
# SLIDE 5 — optimization: shorten outage (node-timeout A/B)
# ============================================================
s5 = prs.slides.add_slide(blank)
heading(s5, "优化：缩短故障转移中断时间（A/B 实测）", "唯一变量 = cluster-node-timeout；检测窗口最大且唯一可观调，是性价比最高的一项")

# left: bar chart outage seconds
cd5 = CategoryChartData()
cd5.categories = ["15000ms\n基线", "5000ms\n推荐", "3000ms\n激进"]
cd5.add_series("客户端写入中断(s)", (20.7, 8.8, 4.8))
g5 = s5.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                         Inches(0.55), Inches(1.4), Inches(6.0), Inches(3.5), cd5)
c5 = g5.chart
c5.has_legend = True
c5.legend.position = XL_LEGEND_POSITION.BOTTOM
c5.legend.include_in_layout = False
c5.legend.font.size = Pt(10)
p5 = c5.plots[0]; p5.gap_width = 70
p5.has_data_labels = True
p5.data_labels.font.size = Pt(12); p5.data_labels.font.bold = True
p5.data_labels.number_format = '0.0"s"'; p5.data_labels.number_format_is_linked = False
c5.series[0].format.fill.solid(); c5.series[0].format.fill.fore_color.rgb = CRIMSON
c5.category_axis.tick_labels.font.size = Pt(10)
c5.value_axis.tick_labels.font.size = Pt(9)

txt(s5, 0.55, 5.0, 6.0, 0.5,
    [[("15000 → 5000ms：中断 ", 12, INK, False), ("降低 58%", 13, GREEN, True),
      ("；→ 3000ms：", 12, INK, False), ("降低 77%", 13, GREEN, True)]])
txt(s5, 0.55, 5.42, 6.0, 0.4,
    [[("中断 ≈ node-timeout + 2~6s（选举+FAIL传播+客户端刷新）", 11, MUTED, False)]])

# right: recommendations (ranked)
tx5 = 6.95
txt(s5, tx5, 1.35, 5.85, 0.4, [[("推荐优化清单（按收益排序）", 15, CRIMSON, True)]])
recs = [
    ("1", "node-timeout 15000→5000ms", "中断 20.7→8.8s。代价：网络抖动更易误判，跨AZ需 RTT≪timeout", DEEP),
    ("2", "客户端快失败 + 主动刷新", "socketTimeout 调小(500~800ms)、命中异常立即刷新 slots", AMBER),
    ("3", "保证多数主存活 + 低滞后副本", "FAIL 需多数主投票；滞后过大副本被禁止升主", BLUE),
    ("4", "计划内切换用 CLUSTER FAILOVER", "副本协商 + offset 对齐，亚秒级且不丢数据", PURPLE),
    ("5", "连接预热 + 连接池", "避免转移后新建连接的额外 RTT 放大尾延迟", GREEN),
]
ry5 = 1.85
for n, t, b, accent in recs:
    rect(s5, tx5, ry5, 5.85, 0.92, WHITE, line=RGBColor(0xD6,0xE2,0xE8), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s5, tx5, ry5, 0.55, 0.92, accent)
    txt(s5, tx5, ry5, 0.55, 0.92, [[(n, 22, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s5, tx5 + 0.72, ry5 + 0.12, 5.0, 0.4, [[(t, 12.5, INK, True)]])
    txt(s5, tx5 + 0.72, ry5 + 0.46, 5.0, 0.42, [[(b, 10, MUTED, False)]], line_spacing=1.0)
    ry5 += 1.02

# bottom synthesis bar
rect(s5, 0.55, 6.25, 6.0, 0.95, RGBColor(0x0E, 0x3D, 0x57), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s5, 0.78, 6.36, 5.6, 0.75,
    [[("工程最优组合", 12.5, RGBColor(0x7F,0xD4,0xF0), True)],
     [("跨 AZ 容灾 + node-timeout=5000 + 客户端快失败重发现 + 幂等重试", 11, ICE, False)]],
    line_spacing=1.1)

out = r"C:\Users\leizha\rediscluster-failover\Redis-Cluster-Failover-汇报.pptx"
prs.save(out)
print("SAVED", out)
