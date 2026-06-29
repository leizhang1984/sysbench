#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""Generate a deployment-architecture deck for J&T (极兔) RocketMQ 4.9.7.

Scenario:
  - 3 NameServers, one in each of AZ 1 / AZ 2 / AZ 3.
  - broker-a/b/c MASTER nodes all in AZ 1.
  - broker-a/b/c SLAVE  nodes all in AZ 2.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- palette (Ocean Gradient) ----
DEEP   = RGBColor(0x06, 0x5A, 0x82)
TEAL   = RGBColor(0x1C, 0x72, 0x93)
MID    = RGBColor(0x21, 0x29, 0x5C)
ICE    = RGBColor(0xE8, 0xF1, 0xF5)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x21, 0x29, 0x2E)
MUTED  = RGBColor(0x6B, 0x7A, 0x85)
GREEN  = RGBColor(0x1B, 0x9E, 0x77)   # master
AMBER  = RGBColor(0xE8, 0x8A, 0x1A)   # slave / warning
RED    = RGBColor(0xC0, 0x39, 0x2B)
AZ1BG  = RGBColor(0xD5, 0xF0, 0xE4)
AZ2BG  = RGBColor(0xFB, 0xEB, 0xD0)
AZ3BG  = RGBColor(0xD6, 0xE9, 0xF5)
CARD   = RGBColor(0xF4, 0xF8, 0xFA)

EMU = 914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width  = Emu(int(SW * EMU))
prs.slide_height = Emu(int(SH * EMU))
blank = prs.slide_layouts[6]


def rect(slide, x, y, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE, line_w=1.0):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    if isinstance(runs[0], tuple):
        runs = [runs]
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line_spacing
        for t, sz, col, bold in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold
            r.font.name = "Calibri"
    return tb


def node(slide, x, y, w, h, title, sub, fill, tcol=WHITE, scol=None):
    rect(slide, x, y, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    scol = scol or tcol
    txt(slide, x + 0.12, y, w - 0.24, h,
        [[(title, 12.5, tcol, True)], [(sub, 9.5, scol, False)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=1, line_spacing=1.0)


# ============================================================
# SLIDE 1 — title
# ============================================================
s1 = prs.slides.add_slide(blank)
rect(s1, 0, 0, SW, SH, MID)
rect(s1, 0, 0, 4.7, SH, DEEP)

txt(s1, 0.55, 0.6, 3.7, 0.5, [[("极兔速递 · J&T Express", 15, ICE, False)]])
txt(s1, 0.55, 1.15, 3.7, 2.4,
    [[("RocketMQ", 32, WHITE, True)],
     [("4.9.7", 32, WHITE, True)],
     [("典型部署架构", 30, WHITE, True)]],
    line_spacing=1.05)
txt(s1, 0.55, 6.4, 3.8, 0.6,
    [[("Master / Slave · 跨可用区部署", 12, ICE, False)]])

# right side summary cards
cx, cw, gap = 5.15, 7.5, 0.0
rows = [
    ("3 台 NameServer", "分别部署于 可用区 1 / 2 / 3，互相独立、无状态", DEEP),
    ("broker-a / b / c — Master", "三个 Master 节点 全部部署在 可用区 1", GREEN),
    ("broker-a / b / c — Slave", "三个 Slave 节点 全部部署在 可用区 2", AMBER),
    ("复制方式", "经典 Master/Slave 主从复制（非 DLedger），Slave 只读", TEAL),
]
ry = 1.25
for title, sub, col in rows:
    rect(s1, cx, ry, cw, 1.18, RGBColor(0x2A, 0x33, 0x6B), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s1, cx, ry, 0.14, 1.18, col)
    txt(s1, cx + 0.35, ry + 0.16, cw - 0.5, 0.9,
        [[(title, 16, WHITE, True)], [(sub, 11.5, ICE, False)]],
        anchor=MSO_ANCHOR.MIDDLE, space_after=3, line_spacing=1.0)
    ry += 1.38

txt(s1, cx, 6.95, cw, 0.4, [[("RocketMQ 4.9.7  ·  6 brokers (3×M + 3×S)  ·  3 NameServer  ·  3 AZ", 10.5, MUTED, False)]])


# ============================================================
# SLIDE 2 — deployment architecture diagram
# ============================================================
s2 = prs.slides.add_slide(blank)
rect(s2, 0, 0, SW, SH, WHITE)
rect(s2, 0, 0, SW, 1.0, DEEP)
txt(s2, 0.5, 0.16, 12.3, 0.7,
    [[("部署架构图 — Master/Slave 跨 AZ 布局", 22, WHITE, True)]],
    anchor=MSO_ANCHOR.MIDDLE)

# client bar
rect(s2, 4.2, 1.18, 4.93, 0.62, MID, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s2, 4.2, 1.18, 4.93, 0.62,
    [[("Producer / Consumer 客户端", 13, WHITE, True)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# three AZ panels
panel_y, panel_h = 2.15, 5.05
col_w, gap = 4.0, 0.33
col_x = [0.5, 0.5 + col_w + gap, 0.5 + 2 * (col_w + gap)]
az_titles = ["可用区 1 (AZ 1)", "可用区 2 (AZ 2)", "可用区 3 (AZ 3)"]
az_bg = [AZ1BG, AZ2BG, AZ3BG]
az_edge = [GREEN, AMBER, TEAL]

for i in range(3):
    rect(s2, col_x[i], panel_y, col_w, panel_h, az_bg[i],
         line=az_edge[i], shape=MSO_SHAPE.ROUNDED_RECTANGLE, line_w=2.0)
    rect(s2, col_x[i], panel_y, col_w, 0.5, az_edge[i])
    txt(s2, col_x[i], panel_y, col_w, 0.5,
        [[(az_titles[i], 14, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# NameServers — one per AZ
ns_y = panel_y + 0.72
ns_labels = ["NameServer ns-1", "NameServer ns-2", "NameServer ns-3"]
for i in range(3):
    node(s2, col_x[i] + 0.35, ns_y, col_w - 0.7, 0.72, ns_labels[i], "无状态 · 路由注册中心", DEEP)

# brokers
brk_y0 = ns_y + 1.05
brk_h, brk_gap = 0.82, 0.28
names = ["broker-a", "broker-b", "broker-c"]
master_x = col_x[0] + 0.35
slave_x = col_x[1] + 0.35
brk_w = col_w - 0.7
for j in range(3):
    y = brk_y0 + j * (brk_h + brk_gap)
    node(s2, master_x, y, brk_w, brk_h, names[j] + "  Master", "可写主节点 (BID=0)", GREEN)
    node(s2, slave_x, y, brk_w, brk_h, names[j] + "  Slave", "只读副本 (BID=1)", AMBER)
    # replication arrow master -> slave
    ax = master_x + brk_w
    rect(s2, ax + 0.02, y + brk_h / 2 - 0.16, slave_x - ax - 0.04, 0.32,
         TEAL, shape=MSO_SHAPE.RIGHT_ARROW)

# AZ3 note (only NameServer there)
txt(s2, col_x[2] + 0.35, brk_y0 + 0.2, col_w - 0.7, 2.0,
    [[("本场景 AZ 3", 12.5, INK, True)],
     [("仅部署 NameServer ns-3，", 11, MUTED, False)],
     [("不承载 broker 节点。", 11, MUTED, False)],
     [("", 6, MUTED, False)],
     [("作用：保证名字服务", 11, MUTED, False)],
     [("跨 3 AZ 高可用。", 11, MUTED, False)]],
    space_after=2, line_spacing=1.05)

# legend
lg_y = panel_y + panel_h + 0.02
txt(s2, 0.5, lg_y - 0.02, 12.4, 0.3,
    [[("图例：", 10.5, INK, True),
      ("■ Master(可写)  ", 10.5, GREEN, True),
      ("■ Slave(只读)  ", 10.5, AMBER, True),
      ("■ NameServer  ", 10.5, DEEP, True),
      ("→ 主从复制(ASYNC/SYNC)；客户端经 NameServer 拉路由后读写 Master。", 10.5, MUTED, False)]])


# ============================================================
# SLIDE 3 — topology details + characteristics & risks
# ============================================================
s3 = prs.slides.add_slide(blank)
rect(s3, 0, 0, SW, SH, WHITE)
rect(s3, 0, 0, SW, 1.0, DEEP)
txt(s3, 0.5, 0.16, 12.3, 0.7,
    [[("拓扑明细与关键风险", 22, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)

# left: node placement table
tx, ty, tw = 0.5, 1.35, 6.1
txt(s3, tx, ty, tw, 0.4, [[("节点分布", 15, DEEP, True)]])
rows = [
    ("组件", "可用区 1", "可用区 2", "可用区 3"),
    ("NameServer", "ns-1", "ns-2", "ns-3"),
    ("broker-a", "Master", "Slave", "—"),
    ("broker-b", "Master", "Slave", "—"),
    ("broker-c", "Master", "Slave", "—"),
]
rh = 0.72
cw0 = 1.9
cwx = (tw - cw0) / 3
yy = ty + 0.5
for ri, row in enumerate(rows):
    head = ri == 0
    for ci, cell in enumerate(row):
        x = tx + (0 if ci == 0 else cw0 + (ci - 1) * cwx)
        w = cw0 if ci == 0 else cwx
        if head:
            fill = DEEP
        elif ci == 0:
            fill = RGBColor(0xEA, 0xF1, 0xF5)
        else:
            fill = CARD
        rect(s3, x, yy, w, rh, fill, line=RGBColor(0xCC, 0xD8, 0xDE))
        tc = WHITE if head else INK
        bold = head or ci == 0
        col = tc
        if not head and cell == "Master":
            col = GREEN; bold = True
        elif not head and cell == "Slave":
            col = AMBER; bold = True
        txt(s3, x, yy, w, rh, [[(cell, 11.5, col, bold)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    yy += rh

txt(s3, tx, yy + 0.15, tw, 1.6,
    [[("特点", 14, DEEP, True)],
     [("• NameServer 跨 3 AZ → 名字服务高可用，任一 AZ 故障仍可路由。", 11, INK, False)],
     [("• Master 与 Slave 分置 AZ1 / AZ2 → 单机/单盘故障可由从节点兜底数据。", 11, INK, False)],
     [("• 经典主从架构，部署简单、读写路径清晰。", 11, INK, False)]],
    space_after=4, line_spacing=1.05)

# right: risk callouts
rx, rw = 7.0, 5.85
txt(s3, rx, ty, rw, 0.4, [[("关键风险与建议", 15, RED, True)]])
risks = [
    ("AZ 1 整体故障 = 全部 Master 不可用", RED,
     "三个 Master 都在 AZ1，AZ1 宕机时写入全部中断；Slave 在 AZ2 只读、不自动升主。"),
    ("4.9.7 经典主从不自动选主", AMBER,
     "原生 Master/Slave 不会自动故障转移，需人工或借助工具切换；恢复时间取决于运维。"),
    ("ASYNC 复制存在 RPO>0 风险", AMBER,
     "若用异步刷盘/异步复制，Master 宕机时未同步到 Slave 的消息可能丢失。"),
    ("建议", GREEN,
     "① Master 分散到多 AZ；② 采用 DLedger(Raft) 实现自动选主；③ 关键 topic 用 SYNC 复制保 RPO=0。"),
]
yy = ty + 0.5
for title, col, body in risks:
    rect(s3, rx, yy, rw, 1.18, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         line=RGBColor(0xDD, 0xE5, 0xEA))
    rect(s3, rx, yy, 0.14, 1.18, col)
    txt(s3, rx + 0.32, yy + 0.13, rw - 0.5, 0.95,
        [[(title, 13, INK, True)], [(body, 10.8, MUTED, False)]],
        anchor=MSO_ANCHOR.MIDDLE, space_after=3, line_spacing=1.02)
    yy += 1.32

OUT = r"C:\Users\leizha\rocketmq-failover\极兔-RocketMQ-4.9.7-部署架构.pptx"
prs.save(OUT)
print("SAVED", OUT)
