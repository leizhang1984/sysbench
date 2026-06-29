#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""Generate a 2-slide analysis deck for the RocketMQ failover round-2 report."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# ---- palette (Ocean Gradient, topic = infra/reliability) ----
DEEP   = RGBColor(0x06, 0x5A, 0x82)   # deep blue (dominant)
TEAL   = RGBColor(0x1C, 0x72, 0x93)   # teal
MID    = RGBColor(0x21, 0x29, 0x5C)   # midnight (dark bg)
ICE    = RGBColor(0xE8, 0xF1, 0xF5)   # very light
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x21, 0x29, 0x2E)
MUTED  = RGBColor(0x6B, 0x7A, 0x85)
GREEN  = RGBColor(0x1B, 0x9E, 0x77)   # good / 0-loss
AMBER  = RGBColor(0xE8, 0x8A, 0x1A)   # warning / SIGSTOP
RED    = RGBColor(0xC0, 0x39, 0x2B)

EMU = 914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width  = Emu(int(SW * EMU))
prs.slide_height = Emu(int(SH * EMU))
blank = prs.slide_layouts[6]


def rect(slide, x, y, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in (tf.margin_left, ):
        pass
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    if isinstance(runs[0], tuple):
        runs = [runs]
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line_spacing
        for t, sz, col, bold in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold
            r.font.name = "Calibri"
    return tb


# ============================================================
# SLIDE 1 — title + key data
# ============================================================
s1 = prs.slides.add_slide(blank)
rect(s1, 0, 0, SW, SH, MID)                       # dark base
rect(s1, 0, 0, 4.55, SH, DEEP)                    # left band

# left band content
txt(s1, 0.55, 0.55, 3.5, 0.5,
    [[("故障转移测试", 16, ICE, False)]])
txt(s1, 0.55, 1.0, 3.55, 2.0,
    [[("RocketMQ 4.9.7", 30, WHITE, True)],
     [("DLedger 故障转移", 30, WHITE, True)],
     [("实测分析", 30, WHITE, True)]],
    line_spacing=1.05)
txt(s1, 0.55, 3.15, 3.55, 2.4,
    [[("方法要点", 13, RGBColor(0x9F,0xC6,0xD8), True)],
     [("• 先查当前 master 位置", 13, ICE, False)],
     [("• 只在该 master 单台 VM 注入", 13, ICE, False)],
     [("• 每次只打一个组，另一组在线", 13, ICE, False)],
     [("• 开启客户端重试 retries=2", 13, ICE, False)],
     [("• 注入→恢复→下一次，共 4 次", 13, ICE, False)]],
    line_spacing=1.15, space_after=6)
txt(s1, 0.55, 6.65, 3.6, 0.4,
    [[("Azure rocketmqnew-rg · 2组×3副本跨3AZ · 1000 msg/s", 9.5, RGBColor(0x8F,0xA9,0xB8), False)]])

# right area title
txt(s1, 4.95, 0.55, 7.9, 0.55,
    [[("核心数据：4 次注入全部零数据丢失（RPO=0）", 21, WHITE, True)]])

# four stat cards
cards = [
    ("Run 1", "kill broker-a master", "≈8s", "仅延迟抬升", "0", "失败", GREEN),
    ("Run 2", "kill broker-b master", "≈10s", "仅延迟抬升", "0", "失败", GREEN),
    ("Run 3", "SIGSTOP broker-a master", "≈37s", "写入塌陷", "84", "失败", AMBER),
    ("Run 4", "SIGSTOP broker-b master", "≈12s", "写入塌陷", "22", "失败", AMBER),
]
cx0, cy0 = 4.95, 1.35
cw, ch = 3.82, 1.62
gx, gy = 0.22, 0.22
for i, (tag, sub, big, biglab, fail, faillab, accent) in enumerate(cards):
    col = i % 2; row = i // 2
    x = cx0 + col * (cw + gx); y = cy0 + row * (ch + gy)
    rect(s1, x, y, cw, ch, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s1, x, y, 0.12, ch, accent)              # accent strip
    txt(s1, x + 0.28, y + 0.14, cw - 0.4, 0.32, [[(tag + "  ", 14, accent, True), (sub, 12, INK, True)]])
    txt(s1, x + 0.28, y + 0.55, 1.9, 0.9,
        [[(big, 34, DEEP, True)], [(biglab + "（中断）", 11, MUTED, False)]], line_spacing=0.95)
    txt(s1, x + 2.35, y + 0.55, 1.3, 0.9,
        [[(fail, 34, accent, True)], [(faillab + "数", 11, MUTED, False)]], line_spacing=0.95)

# bottom insight bar
rect(s1, 4.95, 4.95, cw*2+gx, 1.95, RGBColor(0x0E, 0x3D, 0x57), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s1, 5.25, 5.12, 7.3, 0.4, [[("关键洞察", 14, RGBColor(0x7F,0xD4,0xF0), True)]])
txt(s1, 5.25, 5.5, 7.35, 1.35,
    [[("• ", 13, GREEN, True), ("kill -9 + 重试 + 双组 topic ⇒ 对应用近乎透明：吞吐不掉、0 失败，仅 8–10s 延迟抬升。", 13, ICE, False)],
     [("• ", 13, AMBER, True), ("SIGSTOP（无 RST）即便开重试仍有中断：线程挂满 3s 超时才改投，恢复 12–37s。", 13, ICE, False)],
     [("• ", 13, RGBColor(0x7F,0xD4,0xF0), True), ("服务端选举一致快（≈7–9s），瓶颈在客户端检测/改投；留有在线组+开重试是降中断的关键。", 13, ICE, False)]],
    line_spacing=1.1, space_after=5)

# ============================================================
# SLIDE 2 — charts + comparison analysis
# ============================================================
s2 = prs.slides.add_slide(blank)
rect(s2, 0, 0, SW, SH, ICE)                        # light bg
rect(s2, 0, 0, SW, 1.0, DEEP)                      # header band
txt(s2, 0.55, 0.22, 10.0, 0.6, [[("数据分析：中断时长、失败数与 RPO 核对", 24, WHITE, True)]],
    anchor=MSO_ANCHOR.MIDDLE)
txt(s2, 10.4, 0.22, 2.4, 0.6, [[("kill vs SIGSTOP", 13, RGBColor(0xBF,0xDD,0xEA), True)]],
    align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

runs = ["Run1\nkill-a", "Run2\nkill-b", "Run3\nstop-a", "Run4\nstop-b"]

# chart 1: user-visible interruption seconds
c1 = CategoryChartData(); c1.categories = runs
c1.add_series("用户可见中断 (秒)", (8, 10, 37, 12))
gf1 = s2.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                          Inches(0.55), Inches(1.25), Inches(4.0), Inches(3.0), c1)
ch1 = gf1.chart; ch1.has_legend = False
ch1.value_axis.has_major_gridlines = True
ch1.value_axis.maximum_scale = 40
plot1 = ch1.plots[0]; plot1.has_data_labels = True
plot1.data_labels.number_format = '0"s"'; plot1.data_labels.number_format_is_linked = False
plot1.data_labels.font.size = Pt(10); plot1.data_labels.font.bold = True
ser1 = plot1.series[0]
for idx, pt in enumerate(ser1.points):
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = GREEN if idx < 2 else AMBER
ch1.category_axis.tick_labels.font.size = Pt(9)
ch1.value_axis.tick_labels.font.size = Pt(9)
txt(s2, 0.55, 4.32, 4.0, 0.3, [[("中断时长：kill 几乎无感，SIGSTOP 明显且有方差", 10.5, MUTED, False)]])

# chart 2: failure count
c2 = CategoryChartData(); c2.categories = runs
c2.add_series("故障期失败数", (0, 0, 84, 22))
gf2 = s2.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                          Inches(4.75), Inches(1.25), Inches(4.0), Inches(3.0), c2)
ch2 = gf2.chart; ch2.has_legend = False
ch2.value_axis.maximum_scale = 100
plot2 = ch2.plots[0]; plot2.has_data_labels = True
plot2.data_labels.font.size = Pt(10); plot2.data_labels.font.bold = True
ser2 = plot2.series[0]
for idx, pt in enumerate(ser2.points):
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = GREEN if idx < 2 else AMBER
ch2.category_axis.tick_labels.font.size = Pt(9)
ch2.value_axis.tick_labels.font.size = Pt(9)
txt(s2, 4.75, 4.32, 4.0, 0.3, [[("失败数：kill 经重试改投=0；SIGSTOP 累积少量超时失败", 10.5, MUTED, False)]])

# RPO table (right column)
tx, ty, tw = 8.98, 1.25, 3.8
rect(s2, tx, ty, tw, 3.07, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s2, tx+0.22, ty+0.16, tw-0.4, 0.35, [[("RPO 全量去重核对", 14, DEEP, True)]])
rows = [
    ("测试", "okTotal", "unique", "丢失", True),
    ("Run1", "241,828", "241,828", "0", False),
    ("Run2", "241,382", "241,382", "0", False),
    ("Run3", "205,537", "205,580", "0", False),
    ("Run4", "229,478", "229,483", "0", False),
]
ry = ty + 0.58
col_x = [tx+0.22, tx+1.15, tx+2.18, tx+3.25]
for ri, (a, b, c, d, hdr) in enumerate(rows):
    yy = ry + ri * 0.38
    if hdr:
        rect(s2, tx+0.12, yy-0.02, tw-0.24, 0.36, RGBColor(0xDD,0xEC,0xF2))
    cells = [(a, INK), (b, INK), (c, INK), (d, GREEN)]
    for cxi, (val, colr) in enumerate(cells):
        bold = hdr or cxi == 3
        txt(s2, col_x[cxi], yy, 1.0, 0.34, [[(val, 11, colr if not hdr else DEEP, bold)]],
            anchor=MSO_ANCHOR.MIDDLE)
txt(s2, tx+0.22, ty+2.6, tw-0.4, 0.4,
    [[("unique ≥ okTotal：at-least-once，非丢失；dup 由重试引入", 9.5, MUTED, False)]],
    line_spacing=0.95)

# bottom: takeaways row (3 cards)
by = 4.85
items = [
    ("可用性分水岭", "至少一个组在线 + topic 跨多组 + 开重试，崩溃类故障可被重试即时旁路", GREEN),
    ("SIGSTOP 救不了阻塞", "无 RST 连接被冻结，必须挂满 sendMsgTimeout(3s) 才改投，SLA 按数十秒估", AMBER),
    ("数据零丢失成立", "ASYNC_FLUSH + DLedger 多数派复制；消费端需幂等以容忍少量重复", DEEP),
]
iw = 4.0; ig = 0.22
for i, (h, b, accent) in enumerate(items):
    x = 0.55 + i * (iw + ig)
    rect(s2, x, by, iw, 2.05, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s2, x, by, iw, 0.12, accent)
    txt(s2, x+0.25, by+0.28, iw-0.5, 0.45, [[(h, 15, accent, True)]])
    txt(s2, x+0.25, by+0.85, iw-0.5, 1.1, [[(b, 12.5, INK, False)]], line_spacing=1.12)

# ============================================================
# SLIDE 3 — server-side metrics (DLedger election + NameServer registration)
# ============================================================
s3 = prs.slides.add_slide(blank)
rect(s3, 0, 0, SW, SH, ICE)
rect(s3, 0, 0, SW, 1.0, DEEP)
txt(s3, 0.55, 0.22, 10.5, 0.6,
    [[("服务端指标分析：DLedger 选举 + NameServer 注册", 23, WHITE, True)]],
    anchor=MSO_ANCHOR.MIDDLE)
txt(s3, 10.4, 0.22, 2.4, 0.6, [[("broker 日志 UTC+8", 12, RGBColor(0xBF,0xDD,0xEA), True)]],
    align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

runs3 = ["Run1\nkill-a", "Run2\nkill-b", "Run3\nstop-a", "Run4\nstop-b"]

# chart: server-side election time vs client recovery time
c3 = CategoryChartData(); c3.categories = runs3
c3.add_series("服务端选举耗时 (秒)", (7, 7, 9, 9))
c3.add_series("客户端恢复耗时 (秒)", (8, 10, 37, 12))
gf3 = s3.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                          Inches(0.55), Inches(1.3), Inches(6.0), Inches(3.35), c3)
ch3 = gf3.chart
ch3.has_legend = True
ch3.legend.position = XL_LEGEND_POSITION.BOTTOM
ch3.legend.include_in_layout = False
ch3.legend.font.size = Pt(10)
ch3.value_axis.maximum_scale = 40
plot3 = ch3.plots[0]; plot3.has_data_labels = True
plot3.data_labels.number_format = '0"s"'; plot3.data_labels.number_format_is_linked = False
plot3.data_labels.font.size = Pt(9); plot3.data_labels.font.bold = True
plot3.series[0].format.fill.solid(); plot3.series[0].format.fill.fore_color.rgb = DEEP
plot3.series[1].format.fill.solid(); plot3.series[1].format.fill.fore_color.rgb = AMBER
ch3.category_axis.tick_labels.font.size = Pt(9)
ch3.value_axis.tick_labels.font.size = Pt(9)
txt(s3, 0.55, 4.7, 6.0, 0.55,
    [[("服务端选举一致地快（≈7–9s）；kill 时客户端与之同步，", 11, MUTED, False)],
     [("SIGSTOP 时客户端却滞后到 12–37s ⇒ 瓶颈在客户端而非服务端。", 11, MUTED, False)]],
    line_spacing=1.05)

# server-side election detail table (right)
tx, ty, tw = 6.85, 1.3, 5.95
rect(s3, tx, ty, tw, 3.35, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s3, tx+0.25, ty+0.16, tw-0.5, 0.35, [[("选举与注册明细（实测日志）", 14, DEEP, True)]])
hdr = ("测试", "新 Leader", "term", "切换 cost", "NS 注册")
data = [
    ("Run1", "a-2 (.12)", "14", "18ms", "3台 OK"),
    ("Run2", "b-0 (.13)", "9", "18ms", "3台 OK"),
    ("Run3", "a-2 (.12)", "16", "603ms", "3台 OK"),
    ("Run4", "b-0 (.13)", "11", "603ms", "3台 OK"),
]
colx = [tx+0.25, tx+1.15, tx+2.7, tx+3.45, tx+4.7]
ry = ty + 0.62
rect(s3, tx+0.12, ry-0.04, tw-0.24, 0.36, RGBColor(0xDD,0xEC,0xF2))
for ci, val in enumerate(hdr):
    txt(s3, colx[ci], ry, 1.3, 0.34, [[(val, 10.5, DEEP, True)]], anchor=MSO_ANCHOR.MIDDLE)
for ri, row in enumerate(data):
    yy = ry + 0.42 + ri * 0.42
    accent = GREEN if ri < 2 else AMBER
    for ci, val in enumerate(row):
        bold = ci == 1
        col = accent if ci == 0 else INK
        txt(s3, colx[ci], yy, 1.4, 0.36, [[(val, 10.5, col, bold)]], anchor=MSO_ANCHOR.MIDDLE)
txt(s3, tx+0.25, ty+2.95, tw-0.5, 0.32,
    [[("kill 走 RST（cost=18ms）；SIGSTOP 靠心跳超时检测（cost=603ms），量级相同。", 9.5, MUTED, False)]],
    line_spacing=0.95)

# bottom: 3 server-side insight cards
by = 5.15
items3 = [
    ("选举快且统一", "四次注入服务端均 7–9s 选出新 Leader 并向 3 台 NS 注册；与故障类型几乎无关", DEEP),
    ("瓶颈在客户端", "kill 客户端恢复≈选举时间；SIGSTOP 服务端≈9s 完成、客户端却拖到 12–37s", AMBER),
    ("无脑裂 · NS 不拖后腿", "term 单调+1，落败副本转 Follower；新 Leader 同一秒完成 register broker[0]", GREEN),
]
iw = 4.0; ig = 0.22
for i, (h, b, accent) in enumerate(items3):
    x = 0.55 + i * (iw + ig)
    rect(s3, x, by, iw, 1.85, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s3, x, by, iw, 0.12, accent)
    txt(s3, x+0.25, by+0.26, iw-0.5, 0.4, [[(h, 14.5, accent, True)]])
    txt(s3, x+0.25, by+0.78, iw-0.5, 1.0, [[(b, 11.5, INK, False)]], line_spacing=1.1)

out = r"C:\Users\leizha\rocketmq-failover\RocketMQ-Failover-Round2.pptx"
prs.save(out)
print("SAVED", out)
