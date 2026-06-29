#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""Deployment-architecture deck for the RocketMQ 4.9.7 DLedger (Raft) cluster we built.

Real topology (rocketmqnew-rg):
  - 2 broker groups (broker-a / broker-b), each a 3-node DLedger Raft group.
  - Each group's 3 replicas span AZ 1 / 2 / 3; preferredLeaderId=n1 (AZ2) -> Leader.
  - 3 independent NameServers, one per AZ.
  - listenPort=10911 (client R/W on Leader), DLedger replication port 40911.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DEEP   = RGBColor(0x06, 0x5A, 0x82)
TEAL   = RGBColor(0x1C, 0x72, 0x93)
MID    = RGBColor(0x21, 0x29, 0x5C)
ICE    = RGBColor(0xE8, 0xF1, 0xF5)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x21, 0x29, 0x2E)
MUTED  = RGBColor(0x6B, 0x7A, 0x85)
GREEN  = RGBColor(0x1B, 0x9E, 0x77)   # Leader
AMBER  = RGBColor(0xE8, 0x8A, 0x1A)
RED    = RGBColor(0xC0, 0x39, 0x2B)
FOLL   = RGBColor(0x1C, 0x72, 0x93)   # follower (solid teal, white text)
PURPLE = RGBColor(0x6A, 0x4C, 0x93)   # AZ3 accent
# light, clearly-distinct panel backgrounds per AZ
AZ1BG  = RGBColor(0xDD, 0xEC, 0xF7)   # light blue   (AZ1)
AZ2BG  = RGBColor(0xD5, 0xF0, 0xE4)   # light green  (AZ2 - Leader)
AZ3BG  = RGBColor(0xEC, 0xE2, 0xF5)   # light purple (AZ3)
CARD   = RGBColor(0xF4, 0xF8, 0xFA)

EMU = 914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width  = Emu(int(SW * EMU))
prs.slide_height = Emu(int(SH * EMU))
blank = prs.slide_layouts[6]


def rect(slide, x, y, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE, line_w=1.0):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    if isinstance(runs[0], tuple):
        runs = [runs]
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line_spacing
        for t, sz, col, bold in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold
            r.font.name = "Calibri"
    return tb


def node(slide, x, y, w, h, title, sub, fill, tcol=WHITE, scol=None):
    rect(slide, x, y, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    scol = scol or tcol
    txt(slide, x + 0.1, y, w - 0.2, h,
        [[(title, 11.5, tcol, True)], [(sub, 9, scol, False)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=1, line_spacing=1.0)


# ============================================================
# SLIDE 1 — title
# ============================================================
s1 = prs.slides.add_slide(blank)
rect(s1, 0, 0, SW, SH, MID)
rect(s1, 0, 0, 4.7, SH, DEEP)

txt(s1, 0.55, 0.6, 3.7, 0.5, [[("Azure · rocketmqnew-rg", 14, ICE, False)]])
txt(s1, 0.55, 1.15, 3.8, 2.6,
    [[("RocketMQ 4.9.7", 30, WHITE, True)],
     [("DLedger (Raft)", 30, WHITE, True)],
     [("集群部署架构", 28, WHITE, True)]],
    line_spacing=1.05)
txt(s1, 0.55, 6.35, 3.9, 0.7, [[("多副本自动选主 · 跨 3 可用区", 12, ICE, False)]])

cx, cw = 5.15, 7.5
rows = [
    ("2 组 × 3 副本 = 6 broker", "broker-a / broker-b，每组一个 3 节点 DLedger Raft 组", DEEP),
    ("每组副本跨 3 AZ", "n0/n1/n2 分处 可用区 1 / 2 / 3，任一 AZ 故障仍保多数派", GREEN),
    ("自动选主 (Raft 多数派)", "Leader 失联 → Follower 心跳超时触发选举，无需外部仲裁", TEAL),
    ("3 台 NameServer 跨 AZ", "无状态路由中心，分别部署于 AZ 1 / 2 / 3", AMBER),
]
ry = 1.2
for title, sub, col in rows:
    rect(s1, cx, ry, cw, 1.16, RGBColor(0x2A, 0x33, 0x6B), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s1, cx, ry, 0.14, 1.16, col)
    txt(s1, cx + 0.35, ry + 0.14, cw - 0.5, 0.9,
        [[(title, 15.5, WHITE, True)], [(sub, 11, ICE, False)]],
        anchor=MSO_ANCHOR.MIDDLE, space_after=3, line_spacing=1.0)
    ry += 1.36

txt(s1, cx, 6.95, cw, 0.4,
    [[("flushDiskType=ASYNC_FLUSH · preferredLeaderId=n1 · listenPort 10911 · DLedger 40911", 10, MUTED, False)]])


# ============================================================
# SLIDE 2 — deployment architecture diagram (by AZ)
# ============================================================
s2 = prs.slides.add_slide(blank)
rect(s2, 0, 0, SW, SH, WHITE)
rect(s2, 0, 0, SW, 1.0, DEEP)
txt(s2, 0.5, 0.16, 12.3, 0.7,
    [[("部署架构图 — DLedger 跨 AZ 布局", 22, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)

# client bar
rect(s2, 3.9, 1.12, 5.5, 0.6, MID, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s2, 3.9, 1.12, 5.5, 0.6,
    [[("客户端 (rocketmq-client) · Producer / Consumer", 12.5, WHITE, True)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

panel_y, panel_h = 2.05, 5.15
col_w, gap = 4.0, 0.33
col_x = [0.5, 0.5 + col_w + gap, 0.5 + 2 * (col_w + gap)]
az_titles = ["可用区 1 (AZ 1)", "可用区 2 (AZ 2) — Leader", "可用区 3 (AZ 3)"]
az_bg = [AZ1BG, AZ2BG, AZ3BG]
az_edge = [TEAL, GREEN, PURPLE]

for i in range(3):
    rect(s2, col_x[i], panel_y, col_w, panel_h, az_bg[i],
         line=az_edge[i], shape=MSO_SHAPE.ROUNDED_RECTANGLE, line_w=2.0)
    rect(s2, col_x[i], panel_y, col_w, 0.5, az_edge[i])
    txt(s2, col_x[i], panel_y, col_w, 0.5,
        [[(az_titles[i], 13.5, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# NameServer per AZ
ns_y = panel_y + 0.7
ns = ["NameServer 10.170.0.4", "NameServer 10.170.0.6", "NameServer 10.170.0.5"]
for i in range(3):
    node(s2, col_x[i] + 0.3, ns_y, col_w - 0.6, 0.62, ns[i], ":9876 · 无状态路由", DEEP)

# brokers: group-a row, group-b row; columns = AZ
brk_y_a = ns_y + 0.95
brk_y_b = brk_y_a + 1.18
brk_h = 0.95
brk_w = col_w - 0.6
# data[group][az] = (title, sub, role)
# az index: 0=AZ1(n0 follower) 1=AZ2(n1 leader) 2=AZ3(n2 follower)
groups = [
    ("broker-a", [("a-0 10.170.0.10", "n0 · Follower", FOLL, WHITE),
                  ("a-1 10.170.0.11", "n1 · Leader", GREEN, WHITE),
                  ("a-2 10.170.0.12", "n2 · Follower", FOLL, WHITE)]),
    ("broker-b", [("b-0 10.170.0.13", "n0 · Follower", FOLL, WHITE),
                  ("b-1 10.170.0.14", "n1 · Leader", GREEN, WHITE),
                  ("b-2 10.170.0.15", "n2 · Follower", FOLL, WHITE)]),
]
for gi, (gname, cells) in enumerate(groups):
    yy = brk_y_a if gi == 0 else brk_y_b
    for az in range(3):
        title, sub, fill, tc = cells[az]
        node(s2, col_x[az] + 0.3, yy, brk_w, brk_h, title, sub, fill, tcol=tc,
             scol=ICE)
    # DLedger Raft replication connectors between the 3 replicas (across AZ)
    arr_y = yy + brk_h / 2 - 0.13
    for left_az in (0, 1):
        ax = col_x[left_az] + 0.3 + brk_w           # right edge of left node
        bx = col_x[left_az + 1] + 0.3               # left edge of right node
        rect(s2, ax + 0.03, arr_y, bx - ax - 0.06, 0.26, GREEN,
             shape=MSO_SHAPE.LEFT_RIGHT_ARROW)

# group labels (left edge)
txt(s2, 0.5, brk_y_a + brk_h + 0.04, col_w, 0.25,
    [[("◀ broker-a 组 = 3 节点 DLedger Raft 组（多数派 2/3 复制）", 9, MUTED, True)]])
txt(s2, 0.5, brk_y_b + brk_h + 0.04, col_w, 0.25,
    [[("◀ broker-b 组 = 3 节点 DLedger Raft 组（多数派 2/3 复制）", 9, MUTED, True)]])

# legend
txt(s2, 0.5, panel_y + panel_h + 0.06, 12.4, 0.3,
    [[("图例：", 10.5, INK, True),
      ("■ Leader(可写 :10911)  ", 10.5, GREEN, True),
      ("■ Follower(副本)  ", 10.5, TEAL, True),
      ("■ NameServer  ", 10.5, DEEP, True),
      ("↔ DLedger(Raft) 复制 :40911 — Leader 失联即由多数派自动选主。", 10.5, MUTED, False)]])


# ============================================================
# SLIDE 3 — placement table + failover mechanism
# ============================================================
s3 = prs.slides.add_slide(blank)
rect(s3, 0, 0, SW, SH, WHITE)
rect(s3, 0, 0, SW, 1.0, DEEP)
txt(s3, 0.5, 0.16, 12.3, 0.7,
    [[("节点分布与故障转移机制", 22, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)

# left placement table
tx, ty, tw = 0.5, 1.35, 6.1
txt(s3, tx, ty, tw, 0.4, [[("节点分布 (BID 0 = Leader)", 15, DEEP, True)]])
rows = [
    ("组件", "AZ 1", "AZ 2", "AZ 3"),
    ("NameServer", "10.170.0.4", "10.170.0.6", "10.170.0.5"),
    ("broker-a", "a-0 Follower", "a-1 Leader", "a-2 Follower"),
    ("broker-b", "b-0 Follower", "b-1 Leader", "b-2 Follower"),
]
rh = 0.78
cw0 = 1.55
cwx = (tw - cw0) / 3
yy = ty + 0.5
for ri, row in enumerate(rows):
    head = ri == 0
    for ci, cell in enumerate(row):
        x = tx + (0 if ci == 0 else cw0 + (ci - 1) * cwx)
        w = cw0 if ci == 0 else cwx
        fill = DEEP if head else (RGBColor(0xEA, 0xF1, 0xF5) if ci == 0 else CARD)
        rect(s3, x, yy, w, rh, fill, line=RGBColor(0xCC, 0xD8, 0xDE))
        tc = WHITE if head else INK
        bold = head or ci == 0
        col = tc
        if not head and "Leader" in cell:
            col = GREEN; bold = True
        elif not head and "Follower" in cell:
            col = TEAL
        txt(s3, x + 0.04, yy, w - 0.08, rh, [[(cell, 10.5, col, bold)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    yy += rh

txt(s3, tx, yy + 0.18, tw, 2.2,
    [[("特点", 14, DEEP, True)],
     [("• 每组 3 副本跨 3 AZ → 任一 AZ 整体故障仍保留多数派 (2/3)，可继续选主。", 11, INK, False)],
     [("• Raft 多数派复制 + ASYNC_FLUSH，已提交消息多数派确认。", 11, INK, False)],
     [("• preferredLeaderId=n1：恢复后 Leader 自动回到 AZ2 的 n1 节点。", 11, INK, False)],
     [("• 客户端经 NameServer 拉路由后读写 Leader；Follower 只读跟随。", 11, INK, False)]],
    space_after=4, line_spacing=1.08)

# right: failover mechanism / measured results
rx, rw = 7.0, 5.85
txt(s3, rx, ty, rw, 0.4, [[("自动故障转移机制 (实测)", 15, GREEN, True)]])
cards = [
    ("DLedger 心跳检测", DEEP,
     "Leader→Follower 心跳；失联后 Follower 心跳超时 (~6s) 触发新一轮选举，term+1。"),
    ("服务端选举 ≈ 7–9s", GREEN,
     "新 Leader 选出 (role=LEADER, SYNC_MASTER) 后同一秒向 3 台 NameServer 注册 broker[0]。"),
    ("零数据丢失 RPO = 0", GREEN,
     "四次故障注入 (kill / SIGSTOP) committed 消息全部存活；多数派已确认即不丢。"),
    ("无脑裂 · term 单调递增", TEAL,
     "每次仅产生一个新 Leader，落败副本统一转 Follower/SLAVE，符合 Raft 语义。"),
]
yy = ty + 0.5
for title, col, body in cards:
    rect(s3, rx, yy, rw, 1.18, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         line=RGBColor(0xDD, 0xE5, 0xEA))
    rect(s3, rx, yy, 0.14, 1.18, col)
    txt(s3, rx + 0.32, yy + 0.12, rw - 0.5, 0.95,
        [[(title, 13, INK, True)], [(body, 10.6, MUTED, False)]],
        anchor=MSO_ANCHOR.MIDDLE, space_after=3, line_spacing=1.02)
    yy += 1.32

OUT = r"C:\Users\leizha\rocketmq-failover\RocketMQ-4.9.7-DLedger-集群架构.pptx"
prs.save(OUT)
print("SAVED", OUT)
