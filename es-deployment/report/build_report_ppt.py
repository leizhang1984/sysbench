from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(r"C:\Users\leizha\es-deployment\report")
IMG = ROOT / "images"
OUT = ROOT / "ES-DSv5-vs-DSv6-汇报版-6页-v3.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)


def add_title(slide, text, color=(255, 255, 255), x=0.6, y=0.35, w=12.2, h=0.9, size=34):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*color)
    p.alignment = PP_ALIGN.LEFT


def add_bullets(slide, lines, x, y, w, h, font_size=20, color=(35, 35, 35), line_space=1.15):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, t in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = t
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(*color)
        p.line_spacing = line_space


def card(slide, x, y, w, h, fill=(245, 247, 250), line=(220, 225, 230)):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(*fill)
    shp.line.color.rgb = RGBColor(*line)
    return shp


# Slide 1 - Executive Summary
s1 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s1, (19, 43, 77))
add_title(s1, "Elasticsearch DSv5 vs DSv6\n性能与故障转移汇报", color=(255, 255, 255), size=40, h=1.4)
add_bullets(
    s1,
    [
        "测试范围：esrally geonames 并行压测 + DSv6 自动故障转移验证",
        "环境：Azure germanywestcentral，ES 6.8.1，3 节点集群（master+data）",
        "结论摘要：DSv6 在相同吞吐下显著降低延迟，且具备自动故障转移能力",
        "报告来源：es-dsv5-vs-dsv6-report.md（2026-06-17）",
    ],
    x=0.8,
    y=2.15,
    w=11.8,
    h=3.8,
    font_size=22,
    color=(235, 240, 250),
)

# Slide 2 - Architecture
s_arch = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s_arch, (247, 249, 252))
add_title(s_arch, "测试部署架构（并行压测，互不干扰）", color=(26, 52, 84), size=32)

img_arch = IMG / "architecture.png"
if img_arch.exists():
    s_arch.shapes.add_picture(str(img_arch), Inches(0.35), Inches(0.95), width=Inches(12.6), height=Inches(6.2))

# Slide 3 - Performance Comparison
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s2, (245, 248, 252))
add_title(s2, "压测结果：DSv6 在相同吞吐下延迟全面更优", color=(20, 40, 70), size=32)

card(s2, 0.7, 1.35, 5.35, 5.5, fill=(255, 255, 255))
add_bullets(
    s2,
    [
        "关键数字（服务时间）：",
        "• p50 平均改善约 9%",
        "• p90 平均改善约 17%",
        "• p99 平均改善约 25%~30%",
        "• default 查询 p99 改善 51.4%（20.28ms -> 9.86ms）",
        "",
        "索引路径：",
        "• 累计索引耗时 -1.9%",
        "• 累计 merge 耗时 -8.8%",
        "• 查询吞吐两端一致（固定目标吞吐模型）",
    ],
    x=1.0,
    y=1.7,
    w=4.9,
    h=4.9,
    font_size=17,
    color=(40, 45, 55),
)

img_perf_pct = IMG / "improvement_pct.png"
img_perf_base = IMG / "improvement_baseline.png"
if img_perf_pct.exists():
    s2.shapes.add_picture(str(img_perf_pct), Inches(6.25), Inches(1.45), width=Inches(6.35), height=Inches(2.3))
if img_perf_base.exists():
    s2.shapes.add_picture(str(img_perf_base), Inches(6.25), Inches(3.95), width=Inches(6.35), height=Inches(2.3))

box = s2.shapes.add_textbox(Inches(6.3), Inches(6.25), Inches(6.3), Inches(0.45))
pp = box.text_frame.paragraphs[0]
pp.text = "图：上-相对改善率；下-DSv5/DSv6 p99 绝对值基线"
pp.font.size = Pt(12)
pp.font.color.rgb = RGBColor(95, 105, 120)

# Slide 4 - Failover Architecture (Before/After)
s_fo_arch = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s_fo_arch, (248, 250, 252))
add_title(s_fo_arch, "故障转移架构：故障前后业务连续性与数据持久化", color=(22, 52, 88), size=30)

img_fo_arch = IMG / "failover_architecture_before_after.png"
if img_fo_arch.exists():
    s_fo_arch.shapes.add_picture(str(img_fo_arch), Inches(0.42), Inches(1.15), width=Inches(12.45), height=Inches(5.95))

# Slide 5 - Failover Results
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s3, (250, 250, 250))
add_title(s3, "故障转移验证：支持自动升主，灰色故障检测较慢", color=(18, 46, 40), size=32)

img_fo = IMG / "failover_compare.png"
if img_fo.exists():
    s3.shapes.add_picture(str(img_fo), Inches(0.7), Inches(1.3), width=Inches(7.7), height=Inches(5.3))

card(s3, 8.65, 1.3, 3.95, 5.3, fill=(237, 247, 245), line=(185, 215, 206))
add_bullets(
    s3,
    [
        "场景 A（kill -9）",
        "• ~1s 转 yellow",
        "• +13s 恢复 green",
        "• 失败为 RST 快速失败",
        "",
        "场景 B（SIGSTOP）",
        "• ~93s 后才转 yellow",
        "• +105s 恢复 green",
        "• 失败为 2s 超时挂起",
        "",
        "共同结论",
        "• 自动故障转移成功",
        "• 两场景均零数据丢失",
    ],
    x=8.9,
    y=1.6,
    w=3.5,
    h=4.8,
    font_size=16,
    color=(28, 55, 48),
)

# Slide 6 - Risks & Action Plan
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s4, (255, 255, 255))
add_title(s4, "风险与落地建议（可直接执行）", color=(55, 36, 15), size=32)

card(s4, 0.8, 1.35, 5.9, 5.6, fill=(255, 247, 238), line=(235, 208, 175))
card(s4, 6.95, 1.35, 5.6, 5.6, fill=(241, 248, 255), line=(194, 215, 236))

add_bullets(
    s4,
    [
        "关键风险",
        "1) 灰色故障（无 RST）检测窗口过长",
        "2) 故障转移前客户端易发生超时阻塞",
        "3) 故障转移能力依赖索引副本配置",
        "",
        "当前验证状态",
        "• DSv6：自动故障转移已验证",
        "• failover-test：测试完成并已清理",
    ],
    x=1.1,
    y=1.75,
    w=5.3,
    h=4.9,
    font_size=17,
    color=(90, 62, 26),
)

add_bullets(
    s4,
    [
        "建议行动（优先级从高到低）",
        "1) 生产关键索引设 replicas >= 1",
        "2) 调优 discovery.zen.fd 参数缩短检测窗口",
        "3) 客户端启用拓扑感知 + 快速剔除坏节点",
        "4) 增加节点健康探测与自动告警",
        "",
        "预期收益",
        "• 降低故障恢复时间",
        "• 降低超时对业务的可见影响",
        "• 稳定提升尾延迟与可用性",
    ],
    x=7.25,
    y=1.75,
    w=5.1,
    h=4.9,
    font_size=17,
    color=(30, 56, 90),
)

# Footer on all slides
total = len(prs.slides)
for i, slide in enumerate(prs.slides, start=1):
    ft = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(12.2), Inches(0.25))
    p = ft.text_frame.paragraphs[0]
    p.text = f"Source: es-dsv5-vs-dsv6-report.md | Slide {i}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(120, 126, 136)
    p.alignment = PP_ALIGN.RIGHT

prs.save(str(OUT))
print(f"WROTE: {OUT}")
